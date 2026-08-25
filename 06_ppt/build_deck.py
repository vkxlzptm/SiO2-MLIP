#!/usr/bin/env python3
# ⚠ 주의: v3 이후 사용자가 pptx를 직접 편집했다. 이 스크립트는 그 편집을 담고 있지 않다.
#        재실행하면 사용자 수정분이 사라지므로, 출력 파일명을 _regen 으로 분리해 둔다.
"""템플릿(초록 띠 마스터) 기반 덱 생성 — python-pptx"""
import copy, os
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = lambda n: os.path.join(HERE, "fig", n)

LFONT = "Arial"          # 템플릿 master titleStyle/bodyStyle 의 latin
KFONT = "맑은 고딕"      # 테마 Hangul
GREEN = RGBColor(0x00, 0xC1, 0x1D)
INK = RGBColor(0x00, 0x00, 0x00)
BODY = RGBColor(0x26, 0x26, 0x26)
MUTE = RGBColor(0x7A, 0x7A, 0x7A)
LINE = RGBColor(0xC9, 0xCF, 0xD2)
HDRBG = RGBColor(0x1F, 0x7A, 0x33)
ROWBG = RGBColor(0xF1, 0xF6, 0xF2)

BAND_Y, BAND_H = 7.14, 0.37
TOP, BOT = 1.18, 7.05
M = 0.60
W = 13.3333

prs = Presentation(os.path.join(HERE, "template.pptx"))
LAY = {l.name: l for l in prs.slide_masters[0].slide_layouts}


# ---------- helpers ----------
def _rpr_font(run, size, bold=False, italic=False, color=None):
    f = run.font
    f.name = LFONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    if color is not None:
        f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag, tf_ in (("a:ea", KFONT), ("a:cs", LFONT)):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", tf_)


def _bullet(p, on=True):
    pPr = p._p.get_or_add_pPr()
    for t in ("a:buNone", "a:buChar", "a:buFont", "a:buAutoNum"):
        for e in pPr.findall(qn(t)):
            pPr.remove(e)
    if on:
        pPr.set("marL", str(int(0.26 * 914400)))
        pPr.set("indent", str(-int(0.26 * 914400)))
        bf = pPr.makeelement(qn("a:buFont"), {"typeface": "Wingdings", "charset": "2"})
        bc = pPr.makeelement(qn("a:buChar"), {"char": "ü"})
        pPr.append(bf)
        pPr.append(bc)
    else:
        pPr.set("marL", "0")
        pPr.set("indent", "0")
        pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tb, tf


def para(tf, first, text, size=16, bold=False, italic=False, color=BODY,
         bullet=False, ls=1.3, after=6, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = ls
    p.space_after = Pt(after)
    r = p.add_run()
    r.text = text
    _rpr_font(r, size, bold, italic, color)
    _bullet(p, bullet)
    return p


def bullets(slide, items, x, y, w, h, size=16, after=6, ls=1.3):
    tb, tf = textbox(slide, x, y, w, h)
    for i, t in enumerate(items):
        para(tf, i == 0, t, size=size, bullet=True, after=after, ls=ls)
    return tb


def caption(slide, txt, x, y, w, align=PP_ALIGN.CENTER, h=0.30):
    tb, tf = textbox(slide, x, y, w, h)
    para(tf, True, txt, size=10, italic=True, color=MUTE, after=0, ls=1.3, align=align)
    return tb


def footnote(slide, txt, size=10):
    tb, tf = textbox(slide, 0.50, BAND_Y, 10.10, BAND_H, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, True, txt, size=size, color=INK, after=0, ls=1.3)
    return tb


def title_of(slide, txt):
    ph = slide.shapes.title
    tf = ph.text_frame
    tf.word_wrap = True
    para(tf, True, txt, size=32, bold=False, color=INK, after=0, ls=1.0)
    return ph


def add(name="제목만"):
    return prs.slides.add_slide(LAY[name])


def _set_border(tc, edges=("L", "R", "T", "B"), color=LINE, pt=0.5):
    tcPr = tc.get_or_add_tcPr()
    for e in edges:
        tag = qn("a:ln%s" % e)
        for old in tcPr.findall(tag):
            tcPr.remove(old)
        ln = tcPr.makeelement(tag, {"w": str(int(pt * 12700)), "cap": "flat",
                                    "cmpd": "sng", "algn": "ctr"})
        fill = ln.makeelement(qn("a:solidFill"), {})
        clr = ln.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (color[0], color[4], color[2])})
        fill.append(clr)
        ln.append(fill)
        tcPr.append(ln)


def cell(tc, text, size=13, bold=False, color=BODY, fill=None, align=PP_ALIGN.LEFT):
    tf = tc.text_frame
    tf.word_wrap = True
    tc.margin_left = In(0.08)
    tc.margin_right = In(0.08)
    tc.margin_top = In(0.02)
    tc.margin_bottom = In(0.02)
    tc.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, True, text, size=size, bold=bold, color=color, after=0, ls=1.15, align=align)
    if fill is not None:
        tc.fill.solid()
        tc.fill.fore_color.rgb = fill
    else:
        tc.fill.background()
    _set_border(tc._tc)


def table(slide, data, x, y, w, colw, rowh, size=13, header=False,
          aligns=None, bolds=None, colors=None, fills=None):
    n_r, n_c = len(data), len(data[0])
    gf = slide.shapes.add_table(n_r, n_c, In(x), In(y), In(w), In(rowh * n_r))
    tbl = gf.table
    tbl.first_row = header
    tbl.horz_banding = False
    for j, cw in enumerate(colw):
        tbl.columns[j].width = In(cw)
    for i in range(n_r):
        tbl.rows[i].height = In(rowh)
        for j in range(n_c):
            al = aligns(i, j) if aligns else PP_ALIGN.LEFT
            cell(tbl.cell(i, j), data[i][j], size=size,
                 bold=bolds(i, j) if bolds else False,
                 color=colors(i, j) if colors else BODY,
                 fill=fills(i, j) if fills else None,
                 align=al)
    return gf


# ================= 표지 =================
s = prs.slides[0]
t = s.shapes.title.text_frame
t.word_wrap = True
para(t, True, "범용 MLIP의 a-SiO₂ 적용", size=40, bold=True, color=INK, after=6, ls=1.15)
para(t, False, "SevenNet-Nano 파일럿", size=40, bold=True, color=INK, after=0, ls=1.15)
sub = s.placeholders[4].text_frame
sub.word_wrap = True
para(sub, True, "BKS · 7net-Nano-4.5 · AIMD(PBE) · 실험 비교", size=24, color=BODY, after=18, ls=1.3)
para(sub, False, "이동현    ·    2026-08-14", size=22, color=MUTE, after=0, ls=1.3)
s.shapes.add_picture(os.path.join(HERE, "logo.png"),
                     In(10.75), In(6.22), In(1.85), In(0.77))

# ================= 1. 실공간 구조 =================
s = add()
title_of(s, "실공간 구조 — 결합길이 · 결합각")
bullets(s, [
    "Si–O 1.605 → 1.635 Å — AIMD(PBE) 1.630 Å과 일치 [1]",
    "원인: 7net-Omni(mpa 채널) 증류 → PBE underbinding 상속",
    "→ 국소는 PBE 재현, 중거리 각은 절반만 이동",
], M, 1.22, 5.16, 1.95, after=11)

data = [
    ["", "Si–O", "O–O", "Si–Si", "O–Si–O", "Si–O–Si"],
    ["expt. [2][3]", "1.61", "2.63", "3.08", "109.4–109.7", "140–150"],
    ["BKS", "1.605", "2.605", "3.155", "109.38", "151.89"],
    ["7net-Nano-4.5", "1.635", "2.655", "3.135", "109.41", "145.05"],
    ["AIMD PBE [1]", "1.630", "2.670", "3.030", "109.47", "138.65"],
]
table(s, data, M, 3.50, 5.16, [1.35, 0.62, 0.62, 0.62, 1.15, 0.80], 0.42, size=12,
      aligns=lambda i, j: PP_ALIGN.LEFT if (j == 0 and i > 0) else PP_ALIGN.CENTER,
      bolds=lambda i, j: (i == 0 or i == 3),
      colors=lambda i, j: RGBColor(0xFF, 0xFF, 0xFF) if i == 0 else (GREEN if i == 3 else BODY),
      fills=lambda i, j: HDRBG if i == 0 else (ROWBG if i % 2 == 1 else None))
caption(s, "제1피크 위치 (Å) · 평균 결합각 (°), ρ = 2.20 g/cm³ · 300 K.  "
           "AIMD 참조 [1]는 120원자 · quench 5 ps로 본 연구의 100배 급랭 — 정답으로 두지 않음",
        M, 5.68, 5.16, align=PP_ALIGN.LEFT, h=0.60)

s.shapes.add_picture(FIG("fig_pdf.png"), In(6.05), In(1.30), In(3.39), In(4.44))
s.shapes.add_picture(FIG("fig_bad_vert.png"), In(9.69), In(1.30), In(3.04), In(4.44))
caption(s, "부분 동경분포 g(r)", 6.05, 5.82, 3.39)
caption(s, "결합각 분포", 9.69, 5.82, 3.04)
footnote(s, "[1] Dechant, JPCC 130, 7148 (2026)   ·   [2] Wright, JNCS 179, 84 (1994)   ·   "
            "[3] Dupree & Pettifer, Nature 308, 523 (1984)")

# ================= 2. 역공간 확인 =================
s = add()
title_of(s, "역공간 확인 — FSDP 위치는 불변")
bullets(s, [
    "FSDP 위치 1.593 → 1.590 Å⁻¹ — 두 방법 차이 없음, 실험 대비 +6.7 % [4]",
    "FSDP 진폭 1.462 → 1.360 — 실험 1.357과 일치 (국소 구조 기여)",
    "q_FSDP는 중거리 주기성(2π/q)의 역수 → 망의 위상이 지배",
    "→ 실공간과 독립된 경로에서 같은 결론: 위상은 불변",
], M, 1.25, 12.13, 1.60, after=6)

s.shapes.add_picture(FIG("fig_sq.png"), In(2.08), In(2.95), In(9.17), In(3.80))
caption(s, "중성자 구조인자 S(q) — (a) 전 구간, (b) FSDP 확대   [exp.: [4]]", 2.08, 6.80, 9.17)
footnote(s, "[4] Zeidler, PRL 113, 135501 (2014). FSDP = first sharp diffraction peak, 유리 중거리 질서의 표준 지표.")

# ================= 3. 밀도 · 부피탄성률 =================
s = add()
title_of(s, "밀도 · 부피탄성률 — 비교 조건의 영향")
s.shapes.add_picture(FIG("fig_density.png"), In(1.65), In(TOP), In(6.31), In(2.75))
s.shapes.add_picture(FIG("fig_bulkmod2.png"), In(8.31), In(TOP), In(3.36), In(2.75))
caption(s, "좌: E–V / virial P–V → ρ₀ · K₀          우: 동일 밀도에서 읽은 K(ρ)          [exp.: [5][7]]",
        1.65, 4.02, 10.02)
bullets(s, [
    "밀도 오차 +6.55 % → +0.84 %, 7.8배 감소",
    "K₀(각자 평형): BKS 34.3 / 7net 43.2 GPa — 표면상 BKS 우세",
    "동일 밀도에선 3 % 이내 일치 — K₀′ < 0이라 BKS의 과대 밀도가 K를 낮춤 (오차 상쇄)",
    "ρ = 2.20: 45.3 / 43.9 vs 실험 37 GPa [5] — 둘 다 약 20 % 과대. DFT(GGA) 문헌도 40–46 GPa [6]",
    "→ 오차 출처는 포텐셜이 아니라 구조 생성 프로토콜",
], M, 4.42, 12.13, 2.55, after=6)
footnote(s, "[5] Deschamps, Sci. Rep. 4, 7193 (2014)   ·   [6] Roy & Bongiorno, JPCC 128, 21220 (2024)   ·   "
            "[7] Mazurin et al., Handbook of Glass Data, Part A (1983)")

# ================= 4. 계산 비용 =================
s = add()
title_of(s, "계산 비용 — classical MD 대비 10³배")
bullets(s, [
    "2,160원자 · i5-11600K 6c/12t · GPU 없음 · cutoff 4.5 Å · ρ = 2.20",
    "BKS 1,347,000 vs 7net 1,083 atom-step/s → 1,244배 (5 ps: 2시간 46분 vs 4초)",
    "병렬 천장: OMP 1.83배 / 독립 프로세스 1.98배 → 메모리 대역폭 제한",
    "AIMD 대비: MLIP O(N) vs DFT O(N³) — 본 연구에서 미측정",
], M, TOP + 0.45, 5.05, 4.4, after=20)
s.shapes.add_picture(FIG("fig_speed.png"), In(5.95), In(2.25), In(7.00), In(3.05))
caption(s, "(a) 처리량   (b) 병렬 확장성 — 2,160원자, GPU 없음", 5.95, 5.42, 7.00)
footnote(s, "처리량 = 원자 수 × 스텝 수 ÷ 벽시계 시간(atom-step/s). 배수는 조건 의존 — cutoff 5.5 Å · ρ = 2.607에서는 2,462배.")

# ================= 5. 한계 · next step =================
s = add()
title_of(s, "한계 및 next step")
bullets(s, [
    "7net melt-quench · 어닐링 미수행 → 위상은 BKS 것 (ring 분포 동일, MSD 0.041 Å²)",
    "냉각속도 5×10¹² K/s (실험 대비 ~10¹³배) · 셀 2,160원자 · excess energy 미계산",
    "구조 생성: 용융 중 NPT → 부피 붕괴 (2.26 → 2.62 g/cm³) → NVT 고정 부피로 전환",
    "→ next: 7net melt-quench + 어닐링 [8] — 현 CPU 기준 20 ps ≈ 11시간, GPU 1장이면 실용 범위",
], M, TOP + 0.55, 7.05, 4.4, after=20)
s.shapes.add_picture(FIG("fig_rings.png"), In(8.05), In(1.85), In(4.90), In(3.43))
caption(s, "King ring 크기 분포 — BKS와 7net 이완 구조가 동일", 8.05, 5.42, 4.90)
footnote(s, "[8] Erhard et al., npj Comput. Mater. 8, 90 (2022)")

# ================= 부록 A =================
s = add()
title_of(s, "Appendix. Computational details (1/2)")
rows = [
    ["계", "a-SiO₂ 2,160 atoms (720 Si / 1,440 O), ρ = 2.20 g/cm³, cubic PBC, L ≈ 30.4 Å"],
    ["초기 구조", "Materials Project mp-554089 (Pna2₁, 12 atoms/cell) → 6 × 6 × 5 replicate"],
    ["BKS 포텐셜", "Buckingham + Coulomb, q(Si) = +2.4 e / q(O) = −1.2 e, r_c = 10 Å, PPPM 1×10⁻⁴, 분산항 tail 보정 on"],
    ["단거리 안전항", "r⁻¹² 반발 core (WCA형): O–O σ = 1.50 Å, ε = 3.0 eV / Si–O σ = 1.10 Å, ε = 4.0 eV"],
    ["MLIP", "7net-Nano-4.5 (7net-Omni mpa 채널에서 증류), cutoff 4.5 Å, LAMMPS pair_style e3gnn (TorchScript)"],
    ["구조 생성", "NVT 고정 부피 melt-quench: 300 → 4000 K 50 ps → 4000 K 200 ps 용융 → 4000 → 300 K (100 K 간격 × 20 ps) → 300 K 100 ps"],
    ["냉각속도", "5 × 10¹² K/s"],
    ["용융 검증", "MSD 65.4 Å² (4000 K 200 ps, 시간에 선형) — 확산 확인 후에만 quench 진행"],
    ["적분 · 열욕", "timestep 1 fs, Nosé–Hoover NVT, T_damp = 0.1 ps"],
    ["neighbor", "bin, skin 2.0 Å (BKS) / 1.0 Å (7net), every 1 delay 0 check yes, comm cutoff 9.0 Å (7net)"],
    ["7net MD", "300 K NVT, equilibration 2 ps + production 3 ps, 궤적 20 step 간격 저장"],
    ["300 K NPT 대조", "BKS tail on, 100 ps 시간평균 → ρ = 2.3401 g/cm³ (0 K virial 2.3442와 0.17 % 일치)"],
]
table(s, rows, M, TOP + 0.05, 12.13, [2.30, 9.83], 0.46, size=12,
      bolds=lambda i, j: j == 0,
      colors=lambda i, j: INK if j == 0 else BODY,
      fills=lambda i, j: ROWBG if i % 2 == 0 else None)

# ================= 부록 B =================
s = add()
title_of(s, "Appendix. Computational details (2/2) · References")
rows = [
    ["E–V 스캔", "셀 고정 · 원자만 0 K 이완, 7점. 부피비 f = 0.94–1.06 (7net) / 0.90–1.02 (BKS)"],
    ["최소화 수렴", "etol = 0, ftol = 1×10⁻³ eV/Å, maxiter/maxeval 150 (etol > 0은 셀 자유도에서 조기 종료)"],
    ["상태방정식", "BM3. E(V) 4-파라미터와 virial P(V) 3-파라미터를 독립 피팅, K₀는 P(V) 경로를 채택"],
    ["g(r)", "200 bins, r ≤ 8 Å, 300 K 시간평균. n(Si–O) = 4.001 (컷오프 2.0 Å), Si 4배위 99.86 %"],
    ["S(q)", "주기 셀의 정합 q 벡터 직접 합산 (r 절단 · 창 함수 없음), 껍질 평균 Δq = 2π/L = 0.207 Å⁻¹"],
    ["FSDP 추출", "가우시안 + 선형 배경 피팅 (최댓값은 잡음에 상향 편향). 중성자 가중 0.0694 : 0.3880 : 0.5426"],
    ["ring 통계", "King 기준, 삼중항 (i, j, k) 계수"],
    ["소프트웨어", "LAMMPS stable_2Aug2023_update3 · SevenNet 0.13.0 · PyTorch 2.13.0+cpu · ASE 3.29.0"],
    ["하드웨어", "Intel i5-11600K 6c/12t, RAM 15 GB, GPU 없음"],
]
table(s, rows, M, TOP + 0.02, 12.13, [2.30, 9.83], 0.38, size=12,
      bolds=lambda i, j: j == 0,
      colors=lambda i, j: INK if j == 0 else BODY,
      fills=lambda i, j: ROWBG if i % 2 == 0 else None)

refs = [
    "[1]  G. Dechant, K. Muralidhar, Y. Ma, J. Phys. Chem. C 130, 7148 (2026).",
    "[2]  A. C. Wright, J. Non-Cryst. Solids 179, 84 (1994).",
    "[3]  R. Dupree and R. F. Pettifer, Nature 308, 523 (1984).",
    "[4]  A. Zeidler, K. Wezka, R. F. Rowlands et al., Phys. Rev. Lett. 113, 135501 (2014).",
    "[5]  T. Deschamps, J. Margueritat, C. Martinet et al., Sci. Rep. 4, 7193 (2014).",
    "[6]  U. C. Roy and A. Bongiorno, J. Phys. Chem. C 128, 21220 (2024).",
    "[7]  O. V. Mazurin et al., Handbook of Glass Data, Part A (Elsevier, 1983).",
    "[8]  L. C. Erhard, J. Rohrer, K. Albe, V. L. Deringer, npj Comput. Mater. 8, 90 (2022).",
    "[9]  B. W. H. van Beest, G. J. Kramer, R. A. van Santen, Phys. Rev. Lett. 64, 1955 (1990).",
]
tb, tf = textbox(s, M, 4.72, 12.13, 2.30)
para(tf, True, "References", size=12, bold=True, color=INK, after=3, ls=1.2)
for r in refs:
    para(tf, False, r, size=10, color=BODY, after=0, ls=1.12)

out = os.path.join(HERE, "SiO2_MLIP_SevenNet_pilot_v4_regen.pptx")
prs.save(out)
print("saved:", out, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
