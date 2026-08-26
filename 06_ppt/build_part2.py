#!/usr/bin/env python
"""Part 2 덱 내용 채우기 — v10(Part 1) 테마를 그대로 물려받은 골격에 텍스트·그림을 넣는다.

★ 왜 이런 방식인가
  v10 의 마스터(초록 제목선·하단 각주 띠·페이지 번호·글꼴)를 그대로 쓰려면 새로 만드는
  것보다 **v10 을 복제해 슬라이드를 갈아끼우는 편**이 확실하다. 골격은
  `add_slide.py` 로 만들고(슬라이드 유형 3종을 복제), 이 스크립트는 그 위에
  텍스트와 그림만 교체한다.

  골격 슬라이드 유형 (v10 원본 → Part 2 용도)
    S6형 (2단 텍스트 + 우하단 그림)  -> 2, 3, 8, 9
    S5형 (좌 텍스트 + 우 그림 1개)   -> 4, 6, 7
    S3형 (하단 텍스트 + 그림 2개)    -> 5
    S1(표지) -> 1,   S7(표) -> 10

★ 서식 보존 규칙 (korean-ppt-standards)
  - `text_frame.text = ...` 를 쓰지 않는다. 문단을 통째로 갈아엎으면 run 서식이 사라진다.
  - 대신 **첫 run 의 서식을 복제해** 새 run 을 만든다 (`clone_run`).
  - 한글 run 은 `ea` (동아시아 글꼴) 를 반드시 지정한다. 없으면 기기마다 다르게 렌더된다.
  - 본문 16pt 고정, 제목 32pt. 줄간격 1.3.
  - 줄바꿈 속성(eaLnBrk/latinLnBrk/hangingPunct)은 마지막에 ko_pptx.py apply 로 일괄 적용.
"""
import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
FIG = ROOT / "04_analysis/fig"
SRC = HERE / "_part2_skeleton.pptx"
OUT = HERE / "SiO2_MLIP_SevenNet_pilot_part2_v1.pptx"

NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
BLUE = RGBColor(0x26, 0x7B, 0xB6)     # BKS
PURPLE = RGBColor(0x8E, 0x44, 0xAD)   # 7net on BKS-net
RED = RGBColor(0xD7, 0x29, 0x2A)      # 7net on 7net-net
GRAY = RGBColor(0x59, 0x59, 0x59)


# ------------------------------------------------------------------ helpers
def set_ea(run, face="맑은 고딕"):
    """한글 글꼴(ea) 명시. 비워두면 Windows/macOS/Linux 가 제각각 대체한다."""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(NS + "ea")
    if ea is None:
        ea = rPr.makeelement(NS + "ea", {})
        latin = rPr.find(NS + "latin")
        (latin.addnext(ea) if latin is not None else rPr.append(ea))
    ea.set("typeface", face)


def clone_run(para, proto, text, *, size=None, bold=False, color=None):
    """proto run 의 서식을 복제해 새 run 을 붙인다. 서식 유실 없이 문구만 바꾸는 방법."""
    r = copy.deepcopy(proto._r)
    para._p.append(r)
    from pptx.text.text import _Run
    run = _Run(r, para)
    run.text = text
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    set_ea(run)
    return run


def fill(shape, lines, *, size=16, lead=1.3, gap_pt=11):
    """텍스트 프레임을 lines 로 교체.

    lines 원소:
      str                      -> 기본 서식 한 줄
      (text, dict)             -> dict 로 size/bold/color 지정
      [(text, dict), ...]      -> 한 문단 안 여러 run (부분 강조)
    """
    tf = shape.text_frame
    proto = None
    for p in tf.paragraphs:
        if p.runs:
            proto = p.runs[0]
            break
    if proto is None:
        raise RuntimeError("서식 원본 run 이 없다")
    # 기존 문단 제거 (첫 문단만 남겨 서식 골격 유지)
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]
    # ★ run 만 지우면 <a:br>(shift+enter) 이 남아 문단 앞에 빈 줄이 생긴다.
    #   원본 템플릿 문단에 br 이 섞여 있는 경우가 흔하다 — pPr 만 남기고 전부 지운다.
    for child in list(first._p):
        if not child.tag.endswith("}pPr"):
            first._p.remove(child)

    made = 0
    gap = False
    for item in lines:
        # 빈 항목 = 간격 지시. 빈 문단을 만들면 불릿 기호가 홀로 찍힌다.
        if item is None or (isinstance(item, str) and not item.strip()):
            gap = True
            continue
        para = first if made == 0 else tf.add_paragraph()
        if made > 0:
            pPr = first._p.find(NS + "pPr")
            if pPr is not None:
                para._p.insert(0, copy.deepcopy(pPr))
        segs = item if isinstance(item, list) else [item]
        for seg in segs:
            if isinstance(seg, tuple):
                clone_run(para, proto, seg[0], **seg[1])
            else:
                clone_run(para, proto, seg, size=size)
        para.line_spacing = lead
        if gap and made > 0:
            para.space_before = Pt(gap_pt)
            gap = False
        made += 1

    # run 이 하나도 없는 문단이 남으면 불릿 기호만 홀로 찍힌다 — 반드시 제거한다.
    for para in list(tf.paragraphs):
        if not para.runs:
            para._p.getparent().remove(para._p)
    return shape


def set_title(slide, text):
    """제목 = 상단 20% 안에서 가장 큰 텍스트 도형."""
    cand = [s for s in slide.shapes
            if s.has_text_frame and s.text_frame.text.strip() and s.top < Inches(1.2)]
    cand.sort(key=lambda s: -s.width)
    fill(cand[0], [text], size=32)
    return cand[0]


def set_footnote(slide, text):
    """하단 각주 띠 = 아래 8% 안, 폭 넓은 도형."""
    for s in slide.shapes:
        if s.has_text_frame and s.top > Inches(6.9) and s.width > Inches(6):
            fill(s, [text], size=10.5, lead=1.15)
            return s
    return None


def drop_groups(slide):
    """그림 그룹을 전부 제거하고 그 좌표를 돌려준다 (새 그림을 같은 자리에 놓기 위해)."""
    boxes = []
    for s in list(slide.shapes):
        if s.shape_type == 6 or s.shape_type == 13:
            boxes.append((s.left, s.top, s.width, s.height))
            s._element.getparent().remove(s._element)
    return boxes


def put_pic(slide, png, left, top, width, caption=None, cap_size=10):
    """그림 + 캡션. 가로 기준으로 넣고 세로는 비율로 따라간다."""
    pic = slide.shapes.add_picture(str(png), left, top, width=width)
    if caption:
        tb = slide.shapes.add_textbox(left, top + pic.height + Inches(0.04),
                                      width, Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = caption
        r.font.size = Pt(cap_size)
        r.font.italic = True
        r.font.color.rgb = GRAY
        r.font.name = "Arial"
        set_ea(r)
        p.alignment = 2  # center
        p.line_spacing = 1.15
    return pic


# ------------------------------------------------------------------ build
prs = Presentation(str(SRC))
S = prs.slides

# ---------- 1. 표지 ----------
s = S[0]
tx = [sh for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
tx.sort(key=lambda sh: sh.top)
fill(tx[0], [[("Universal MLIP", {"size": 40, "bold": True}),
              ("의 a-SiO", {"size": 40, "bold": True}),
              ("₂ ", {"size": 40, "bold": True}),
              ("적용", {"size": 40, "bold": True})],
             [("7net 자체 melt-quench", {"size": 40, "bold": True})]], size=40)
fill(tx[1], [("BKS · 7net-Nano-4.5 · AIMD(PBE) · 실험 비교", {"size": 24}),
             ("Part 2", {"size": 24, "bold": True}),
             ("2026. 08. 27", {"size": 20}),
             ("이동현", {"size": 20})], size=24, lead=1.35)

# ---------- 2. Part 1 이 남긴 질문 ----------
s = S[1]
set_title(s, "Part 1 이 남긴 질문 — 위상은 정말 못 고치나")
drop_groups(s)
body = sorted([sh for sh in s.shapes if sh.has_text_frame and Inches(1.0) < sh.top < Inches(6.5)],
              key=lambda sh: sh.left)
fill(body[0], [
    ("[ Part 1 의 결론 ]", {"size": 16, "bold": True}),
    "7net = BKS 가 만든 망을 물려받아 이완만 수행",
    [("→ 국소는 고침", {"size": 16}), (" · ", {"size": 16}), ("위상은 못 고침", {"size": 16, "bold": True})],
    "",
    ("[ 근거 2 가지 ]", {"size": 16, "bold": True}),
    [("① ring 분포가 BKS 와 ", {"size": 16}), ("동일", {"size": 16, "bold": True})],
    [("② FSDP 위치 1.593 → 1.590 Å⁻¹ ", {"size": 16}), ("(변화 없음)", {"size": 16, "bold": True})],
    "   실험 1.492 대비 둘 다 +6.7 %",
])
fill(body[1], [
    ("[ 그래서 Part 2 의 질문 ]", {"size": 16, "bold": True}),
    "",
    [("7net 이 ", {"size": 16}), ("스스로 망을 만들면", {"size": 16, "bold": True, "color": RED}),
     ("\n위상이 바뀌는가?", {"size": 16, "bold": True})],
    "",
    ("[ 답하는 방법 ]", {"size": 16, "bold": True}),
    "7net 자체 melt-quench 수행",
    "= Part 1 의 Next step 항목 그대로",
    "",
    "★ 냉각률을 맞춘 BKS 통제런을 짝지어",
    "   포텐셜 효과 / 위상 효과를 분리",
])
set_footnote(s, "Part 1 = SiO2_MLIP_SevenNet_pilot_v10.  FSDP = first sharp diffraction peak, 유리 중거리 질서의 표준 지표.")

# ---------- 3. 계산 설계 ----------
s = S[2]
set_title(s, "계산 설계 — 통제군이 핵심")
drop_groups(s)
body = sorted([sh for sh in s.shapes if sh.has_text_frame and Inches(1.0) < sh.top < Inches(6.5)],
              key=lambda sh: sh.left)
fill(body[0], [
    ("[ 7net 자체 melt-quench ]", {"size": 16, "bold": True}),
    "BKS 4000 K 액체 → 재평형(MSD gate)",
    [("→ ", {"size": 16}), ("2×10¹³ K/s", {"size": 16, "bold": True}), (" 급랭 → 300 K", {"size": 16})],
    "셀 2,160원자 · ρ = 2.20 g/cm³ 고정 NVT",
    "",
    ("[ 셀·밀도를 고정한 이유 ]", {"size": 16, "bold": True}),
    "ring · S(q) · 밀도가 모두 셀 크기에 민감",
    "→ 바꾸면 Part 1 결과와 비교 불가",
    "",
    [("비용: 7net 본런 ", {"size": 16}), ("122 시간", {"size": 16, "bold": True, "color": RED}),
     (" (시드 1개가 한계의 근원)", {"size": 16})],
])
fill(body[1], [
    ("[ 통제군 — 이게 설계의 핵심 ]", {"size": 16, "bold": True}),
    "",
    [("① 냉각률 매칭 BKS 런", {"size": 16, "bold": True})],
    "   같은 2×10¹³ K/s · 독립 시드 3개",
    "   → 포텐셜 효과 vs 위상 효과 분리",
    "",
    [("② 냉각률 스윕 BKS", {"size": 16, "bold": True})],
    "   1×10¹² ~ 1×10¹⁴ K/s (7 rates)",
    "   → 관측된 차이가 냉각률 탓인지 판정",
    "",
    "★ 통제군 없이는 어떤 차이도 해석 불가",
])
set_footnote(s, "MSD gate = 4000 K 액체가 확산 영역에 도달했는지 확인 후 급랭 시작.  급랭 속도는 100 K / 5 ps 단위로 단계 냉각.")

# ---------- 4. Tg ----------
s = S[3]
set_title(s, "유리전이 — 7net 의 액체는 다른 액체다")
boxes = drop_groups(s)
body = [sh for sh in s.shapes if sh.has_text_frame and Inches(1.0) < sh.top < Inches(6.9)
        and sh.width < Inches(8)]
fill(body[0], [
    [("kinetic-arrest T", {"size": 16, "bold": True}), ("g", {"size": 12, "bold": True}),
     (" : 7net ", {"size": 16, "bold": True}), ("2000 K", {"size": 16, "bold": True, "color": RED}),
     (" vs BKS ", {"size": 16, "bold": True}), ("2800 K", {"size": 16, "bold": True, "color": BLUE}),
     ("  → 800 K 낮음", {"size": 16, "bold": True})],
    "냉각률 매칭 · BKS 3시드 중앙값 (실측 산포 2500~3000 K)",
    None,
    [("환산", {"size": 16, "bold": True}), (" : BKS T", {"size": 16}), ("g", {"size": 12}),
     (" = 180 ± 40 K / decade → 800 K = ", {"size": 16}),
     ("10⁴·⁴ 배 느린 냉각", {"size": 16, "bold": True})],
    "선형 외삽 · 느린 쪽에서 완만해짐 → 하한값",
    None,
    [("배열 에너지", {"size": 16, "bold": True}), (" : 7net ", {"size": 16}),
     ("249", {"size": 16, "bold": True, "color": RED}),
     (" vs BKS 71–106 meV/atom (", {"size": 16}),
     ("2.4–3.5 배", {"size": 16, "bold": True}), (")", {"size": 16})],
    None,
    [("[ 검증 ] ", {"size": 16, "bold": True}),
     ("유리가지 dE/dT 가 11개 런 전부 1.5k", {"size": 16}), ("B", {"size": 12}),
     ("의 1.04배", {"size": 16})],
    "→ Tg 아래선 두 포텐셜이 구별 안 됨 = 순수 진동",
    "→ 그 기준선 위로 벌어진 양만이 망 재배열의 몫",
])
put_pic(s, FIG / "fig_tg_s4.png", Inches(7.30), Inches(1.25), Inches(5.65),
        "(a) 배열 에너지  (b) MSD 기반 arrest  (c) 냉각률 의존성")
set_footnote(s, "arrest 임계값 ΔMSD = 0.10 Å²/100 K · K = 5 연속.  2×10¹⁴ K/s 는 동적범위 부족으로 제외.  decade = 로그축 10배.")

# ---------- 5. 망 위상 (ring + angle) ----------
s = S[4]
set_title(s, "망 위상 — 고리 분포와 결합각")
drop_groups(s)
# 그림 2개를 위쪽에 나란히 두므로 본문 박스를 아래로 내리고 폭을 넓힌다.
body = [sh for sh in s.shapes if sh.has_text_frame and sh.top > Inches(3.5)
        and sh.top < Inches(6.9) and sh.width > Inches(8)][0]
body.left, body.top = Inches(0.72), Inches(5.28)
body.width, body.height = Inches(11.95), Inches(1.80)
fill(body, [
    [("3-ring", {"size": 16, "bold": True}), (" 1.80 % → ", {"size": 16}),
     ("7.17 %", {"size": 16, "bold": True, "color": RED}), (" (4배, p = 0.038)", {"size": 16}),
     ("      5-ring", {"size": 16, "bold": True}), (" −6.87 %p (", {"size": 16}),
     ("p = 0.021, 가장 유의", {"size": 16, "bold": True, "color": RED}), (")", {"size": 16})],
    [("냉각률 효과 < 실현간 산포", {"size": 16, "bold": True}),
     (" — 5e12 0.4 sd · 5e13 0.9 sd → 냉각률 아티팩트 아님 (3시드 실측)", {"size": 16})],
    [("⚠ 3-ring 은 원인이 아니라 ", {"size": 16, "bold": True}),
     ("증상", {"size": 16, "bold": True, "color": RED}),
     (" — 각 이동 −10.48° 중 3-ring 몫은 ", {"size": 16}),
     ("11 %", {"size": 16, "bold": True}), (", 89 %는 전역적 좁아짐", {"size": 16})],
    [("각 사다리", {"size": 16, "bold": True}),
     (" : BKS 151.9° → 7net(BKS망) 145.2° → ", {"size": 16}),
     ("7net(자기망) 140.7°", {"size": 16, "bold": True, "color": RED}),
     (" → AIMD 138.7°", {"size": 16})],
], gap_pt=4)
put_pic(s, FIG / "fig_rings_s4.png", Inches(0.72), Inches(1.16), Inches(5.20),
        "Ring 분포 — 오차막대는 매칭 냉각률 3시드의 실측 산포")
put_pic(s, FIG / "fig_angle_s4.png", Inches(6.60), Inches(1.16), Inches(5.20),
        "Si–O–Si 각 — 비3원환 산소도 함께 내려간다")
set_footnote(s, "p = 양측 t 검정(2 dof).  표본 3개라 sd 자체가 불확실 → σ 대신 p 를 인용.  7net 은 시드 1개이므로 산포 성분은 전부 BKS 에서 추정.")

# ---------- 6. S(q) ----------
s = S[5]
set_title(s, "S(q) — 위치는 망이 98 % 정한다")
drop_groups(s)
body = [sh for sh in s.shapes if sh.has_text_frame and Inches(1.0) < sh.top < Inches(6.9)
        and sh.width < Inches(8)]
fill(body[0], [
    [("q", {"size": 16, "bold": True}), ("FSDP", {"size": 12, "bold": True}),
     (" : 1.590 (+6.6 %) → ", {"size": 16, "bold": True}),
     ("1.470 ± 0.060 (−1.5 %)", {"size": 16, "bold": True, "color": RED})],
    "d = 2π/q : 3.95 → 4.27 Å   (실험 4.21 Å)",
    [("냉각률 통제런은 1.588 로 그대로 → ", {"size": 16}),
     ("냉각률 아님", {"size": 16, "bold": True})],
    None,
    ("[ ρ = 2.2000 통제 → 분해 가능 ]", {"size": 16, "bold": True}),
    [("위치 : [망] ", {"size": 16}), ("98 %", {"size": 16, "bold": True, "color": RED}),
     (" / [포텐셜] 2 %   → 48 배", {"size": 16})],
    [("진폭 : [망] 70 % / [포텐셜] ", {"size": 16}), ("30 %", {"size": 16, "bold": True}),
     ("   → 국소도 기여", {"size": 16})],
    None,
    [("⚠ 진폭은 무너졌다", {"size": 16, "bold": True, "color": RED}),
     (" : 1.360 → ", {"size": 16}), ("1.123 (−17.2 %)", {"size": 16, "bold": True})],
    "Si–O–Si 각 σ 13.55 → 15.45 로 넓어진 것과 정합",
    "→ \"위치가 맞았다\"만 떼어 인용 금지",
])
put_pic(s, FIG / "fig_sq_s4.png", Inches(7.00), Inches(2.10), Inches(5.95),
        "중성자 S(q) — 색 = 포텐셜+망 조합. (a) 전 구간 (b) FSDP 확대")
set_footnote(s, "[4] A. Zeidler et al., PRL 113, 135501 (2014), digitized.   q_FSDP = 가우시안+선형배경 피팅값.  자기망은 피크가 평평해 불확도 3~4배 — ±0.06 병기.")

# ---------- 7. 밀도·탄성률 ----------
s = S[6]
set_title(s, "밀도 · 부피탄성률 — Part 1 예측의 검증")
drop_groups(s)
body = [sh for sh in s.shapes if sh.has_text_frame and Inches(1.0) < sh.top < Inches(6.9)
        and sh.width < Inches(8)]
fill(body[0], [
    ("[ Part 1 의 예측 ]", {"size": 16, "bold": True}),
    "\"K₀ 오차의 출처는 포텐셜이 아니라 구조 생성 프로토콜\"",
    None,
    [("K @ ρ=2.20 : 43.9 → ", {"size": 16, "bold": True}),
     ("37.7 GPa", {"size": 16, "bold": True, "color": RED}),
     ("   (+18.7 % → ", {"size": 16, "bold": True}),
     ("+2.0 %", {"size": 16, "bold": True, "color": RED}), (")", {"size": 16, "bold": True})],
    "→ 예측 적중",
    None,
    ("★ 단 후반부는 틀렸다", {"size": 16, "bold": True, "color": RED}),
    "\"현실적인 냉각속도까지 내려가야\" 한다고 했으나",
    "S4 는 4배 빠른 냉각인데도 K 가 실험에 붙었다",
    "→ 냉각속도가 아니라 \"누가 망을 만들었나\"가 지배",
    None,
    [("⚠ 밀도는 반대로 나빠졌다", {"size": 16, "bold": True, "color": RED}),
     (" : ρ₀ +0.84 % → ", {"size": 16}), ("−1.84 %", {"size": 16, "bold": True})],
])
put_pic(s, FIG / "fig_bulkmod2.png", Inches(7.55), Inches(1.55), Inches(5.35),
        "K(ρ) — 같은 밀도에서 비교. 별 = 실험 fused silica")
set_footnote(s, "두 곡선의 낙차(6.2 GPa)를 '위상 효과'로 인용하지 말 것 — ev220 스캔은 7점 중 6점이 minimize 미수렴이고 S4 자기망은 전점 수렴이라 같은 자가 아니다.")

# ---------- 8. 종합 ----------
s = S[7]
set_title(s, "종합 — 한쪽을 고치면 다른 쪽이 상한다")
drop_groups(s)
body = sorted([sh for sh in s.shapes if sh.has_text_frame and Inches(1.0) < sh.top < Inches(6.5)],
              key=lambda sh: sh.left)
fill(body[0], [
    ("[ 독립된 두 경로, 같은 패턴 ]", {"size": 16, "bold": True}),
    "",
    [("탄성 (E–V)", {"size": 16, "bold": True})],
    [("   좋아짐 : K@ρ_exp +18.7 → ", {"size": 16}), ("+2.0 %", {"size": 16, "bold": True, "color": RED})],
    [("   나빠짐 : ρ₀ +0.84 → ", {"size": 16}), ("−1.84 %", {"size": 16, "bold": True, "color": BLUE})],
    "",
    [("산란 (S(q))", {"size": 16, "bold": True})],
    [("   좋아짐 : q_FSDP +6.6 → ", {"size": 16}), ("−1.5 %", {"size": 16, "bold": True, "color": RED})],
    [("   나빠짐 : S(q_FSDP) +0.2 → ", {"size": 16}), ("−17.2 %", {"size": 16, "bold": True, "color": BLUE})],
    "",
    "→ 서로 독립인 두 측정에서 같은 구조 → 우연 아님",
])
fill(body[1], [
    ("[ 여섯 지표가 같은 방향 ]", {"size": 16, "bold": True}),
    "3-ring 4배↑ · 배열에너지 2.4–3.5배↑ · Tg 800 K↓",
    "평균 Si–O–Si −10.5° · K 6.2 GPa↓ · ρ₀ 2.7 %↓",
    "",
    [("→ 7net 이 스스로 만든 망은", {"size": 16}),
     ("\n   더 변형되고 · 더 성기고 · 더 무른 망", {"size": 16, "bold": True, "color": RED})],
    "",
    ("[ 한 줄로 ]", {"size": 16, "bold": True}),
    "포텐셜로는 위상을 못 고치고,",
    "망을 새로 만들면 위상은 고쳐지되 다른 쪽이 상한다",
    "",
    [("⚠ 정합성이지 ", {"size": 16, "bold": True}),
     ("인과 사슬의 분리 검증은 아니다", {"size": 16, "bold": True, "color": RED})],
])
set_footnote(s, "Part 1 의 결론 \"국소는 고치고 위상은 못 고친다\"는 BKS 위상을 물려받은 조건에서만 참인 문장이었다.")

# ---------- 9. Limitations & Next step ----------
s = S[8]
set_title(s, "Limitations & Next step")
drop_groups(s)
body = sorted([sh for sh in s.shapes if sh.has_text_frame and Inches(1.0) < sh.top < Inches(6.5)],
              key=lambda sh: sh.left)
fill(body[0], [
    ("[ Limitations ]", {"size": 16, "bold": True}),
    [("7net 계열은 전부 ", {"size": 16}), ("시드 1개", {"size": 16, "bold": True, "color": RED}),
     (" (122 h / 런)", {"size": 16})],
    "→ 모든 산포 성분은 BKS 에서 추정한 값",
    "",
    "ring 분포에는 실험 계열이 없음",
    "→ \"실물 유리와 다르다\"의 근거로 쓸 수 없음",
    "   (실험과 직접 대조되는 것은 S(q) 뿐)",
    "",
    "E–V 의 [냉각률]/[위상] 분해는 폐기",
    "→ 스캔마다 minimize 수렴도가 달라 같은 자가 아님",
    "   밀도가 통제된 S(q) 쪽만 분해 성립",
    "",
    "AIMD 참조 1편 (120원자 · 5 ps quench)",
])
fill(body[1], [
    ("[ Next step ]", {"size": 16, "bold": True}),
    "",
    [("7net 반복 시드", {"size": 16, "bold": True}), (" — 산포를 자체 측정", {"size": 16})],
    "   현 CPU 기준 1런 122시간 → GPU 필요",
    "",
    [("어닐링", {"size": 16, "bold": True}), (" — 급랭 이력 제거 후 재측정", {"size": 16})],
    "",
    [("열팽창계수 α", {"size": 16, "bold": True}), (" — 결정(α-quartz) 먼저 검증", {"size": 16})],
    "",
    [("GPU 벤치", {"size": 16, "bold": True}), (" — 계가 클수록 유리", {"size": 16})],
    "",
    "★ 우선순위는 7net 반복 시드",
    "   지금 모든 유의성 판정이 BKS 산포에 의존",
])
set_footnote(s, "[9] L. C. Erhard et al., npj Comput. Mater. 8, 90 (2022).   급랭 속도 5×10¹² ~ 2×10¹³ K/s 는 실험 대비 ~10¹³배 빠르다.")

# ---------- 10. Appendix ----------
s = S[9]
set_title(s, "Appendix. S4 계산 조건")
tbl = None
for sh in s.shapes:
    if sh.has_table:
        tbl = sh.table
        break
rows = [
    ("계", "a-SiO₂ 2,160 atoms (Si 720 / O 1,440), ρ = 2.20 g/cm³ 고정 NVT"),
    ("출발 구조", "BKS 4000 K 액체 (200 ps 용융) → MSD gate 로 확산 확인 후 급랭"),
    ("7net melt-quench", "2×10¹³ K/s (100 K / 5 ps 단계 냉각) → 300 K, wall time 122 h"),
    ("BKS 통제런", "매칭 2×10¹³ K/s 독립 시드 3개 (90210/90211/90212) + 스윕 1e12~1e14"),
    ("MLIP", "7net-Nano-4.5 (7net-Omni mpa 증류, PBE(+U) 계열), cutoff 4.5 Å"),
    ("Tg 판정", "kinetic arrest: ΔMSD/100 K < 0.10 Å², K = 5 연속.  동적범위 게이트 적용"),
    ("Ring 통계", "King 기준, distinct fraction.  RCUT 1.85~2.15 Å 스윕에서 순위 불변"),
    ("Si–O–Si 각", "prod_*.data 300 K 스냅샷.  트라젝토리 평균과 0.15° 이내 일치 확인"),
    ("E–V 스캔", "0 K, 셀 고정·원자만 이완.  pilot 로 P=0 괄호 → V₀ 대칭 격자 7점"),
    ("최소화 수렴", "판정 = 'linesearch alpha is zero' + 최종 힘 2-norm 0.02~0.05"),
    ("S(q)", "주기 셀 정합 q 벡터 직접 합산 (절단·창 없음), 껍질폭 0.10, stride 5"),
    ("FSDP 추출", "가우시안 + 선형 배경 피팅.  최댓값은 위로 편향되어 사용 금지"),
    ("소프트웨어", "LAMMPS stable_2Aug2023 + SevenNet, Python 3 / numpy / scipy"),
    ("하드웨어", "Intel i5-11600K 6c12t, RAM 15 GB, GPU 없음"),
]
if tbl is not None:
    n = min(len(rows), len(tbl.rows))
    for i in range(n):
        for j, txt in enumerate(rows[i][:len(tbl.columns)]):
            cell = tbl.cell(i, j)
            tf = cell.text_frame
            proto = None
            for p in tf.paragraphs:
                if p.runs:
                    proto = p.runs[0]
                    break
            if proto is None:
                continue
            for p in list(tf.paragraphs)[1:]:
                p._p.getparent().remove(p._p)
            first = tf.paragraphs[0]
            for r in list(first.runs)[1:]:
                r._r.getparent().remove(r._r)
            first.runs[0].text = txt
            set_ea(first.runs[0])
    # 남는 행은 비운다
    for i in range(n, len(tbl.rows)):
        for j in range(len(tbl.columns)):
            tf = tbl.cell(i, j).text_frame
            for p in tf.paragraphs:
                for r in list(p.runs):
                    r.text = ""
set_footnote(s, "[10] B. W. H. van Beest, G. J. Kramer, R. A. van Santen, PRL 64, 1955 (1990).   [1] G. Dechant et al., JPCC 130, 7148 (2026).")

prs.save(str(OUT))
print(f"-> {OUT}")
